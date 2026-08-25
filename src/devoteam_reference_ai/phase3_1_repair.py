from __future__ import annotations

import csv
import hashlib
import io
import json
import os
from collections import Counter, defaultdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable

import fitz
import pandas as pd
import yaml
from PIL import Image, ImageFilter, ImageOps, ImageSequence

from .phase2_utils import sha256_file
from .phase3_pipeline import verify_phase3
from .phase3_quality import image_ink_ratio, normalize_extracted_text, text_metrics
from .phase3_security import redact_text


KEY_COLUMNS = ["document_id", "page_number_1_based"]
FINAL_TEXT_COLUMNS = [
    "extraction_method",
    "ocr_engine",
    "ocr_languages",
    "ocr_confidence",
    "ink_ratio",
    "character_count",
    "word_count",
    "line_count",
    "latin_character_count",
    "arabic_character_count",
    "is_blank",
    "qa_status",
    "qa_reasons",
    "redaction_counts",
    "text_raw",
    "text_redacted",
]


def _now() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")


def _text_sha256(value: str) -> str:
    return hashlib.sha256(value.encode("utf-8")).hexdigest()


def _atomic_json(path: Path, value: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_suffix(path.suffix + ".tmp")
    temporary.write_text(
        json.dumps(value, ensure_ascii=False, indent=2, sort_keys=True),
        encoding="utf-8",
    )
    os.replace(temporary, path)


def load_phase3_1_config(path: Path) -> dict:
    config = yaml.safe_load(path.read_text(encoding="utf-8"))
    if float(config.get("phase", 0)) != 3.1:
        raise ValueError("Expected Phase 3.1 configuration")
    if config.get("security", {}).get("external_llm_enabled"):
        raise ValueError("External LLM calls are forbidden during Phase 3.1")
    return config


def _classification(metrics: dict, confidence: float | None, config: dict) -> tuple[str, list[str]]:
    quality = config["quality"]
    characters = int(metrics["character_count"])
    words = int(metrics["word_count"])
    reasons: list[str] = []
    if characters < int(quality["review_min_characters"]):
        return "FAILED", ["very_low_text_volume"]
    if characters < int(quality["pass_min_characters"]) or words < int(quality["pass_min_words"]):
        reasons.append("low_text_volume")
        return "REVIEW", reasons
    numeric_confidence = float(confidence or 0.0)
    if numeric_confidence < float(quality["ocr_review_min_confidence"]):
        return "FAILED", ["very_low_ocr_confidence"]
    if numeric_confidence < float(quality["ocr_pass_min_confidence"]):
        return "REVIEW", ["low_ocr_confidence"]
    return "PASS", []


def _attempt_score(attempt: dict) -> tuple:
    rank = {"PASS": 3, "REVIEW": 2, "FAILED": 1, "BLANK": 0}.get(
        attempt["qa_status"], 0
    )
    return (
        rank,
        min(int(attempt.get("word_count") or 0), 1000),
        min(int(attempt.get("character_count") or 0), 5000),
        round(float(attempt.get("ocr_confidence") or 0.0), 4),
        -int(attempt.get("attempt_order") or 0),
    )


def _prepare_image(image: Image.Image, variant: str, rotation: int = 0) -> Image.Image:
    prepared = ImageOps.exif_transpose(image.copy()).convert("RGB")
    if rotation:
        prepared = prepared.rotate(rotation, expand=True, fillcolor="white")
    gray = ImageOps.autocontrast(ImageOps.grayscale(prepared))
    # Scans in the source corpus are often a small gray document centered on a
    # large white PDF canvas. Tesseract can return no words for that layout.
    # Crop only near-white outer canvas while keeping a conservative margin.
    content_mask = gray.point(lambda pixel: 255 if pixel < 245 else 0, mode="1")
    bounding_box = content_mask.getbbox()
    if bounding_box:
        left, top, right, bottom = bounding_box
        content_area = max(right - left, 1) * max(bottom - top, 1)
        image_area = max(gray.width * gray.height, 1)
        if 0.05 <= content_area / image_area < 0.98:
            padding_x = max(int((right - left) * 0.015), 8)
            padding_y = max(int((bottom - top) * 0.015), 8)
            gray = gray.crop(
                (
                    max(left - padding_x, 0),
                    max(top - padding_y, 0),
                    min(right + padding_x, gray.width),
                    min(bottom + padding_y, gray.height),
                )
            )
            # Re-stretch contrast after removing the dominant white canvas;
            # otherwise faint gray document text can remain below OCR contrast.
            # The photographed sheet itself is gray; clipping the brightest
            # 10% maps that paper background back to white while retaining ink.
            gray = ImageOps.autocontrast(gray, cutoff=(1, 10))
    if variant == "threshold":
        return gray.point(lambda pixel: 255 if pixel > 185 else 0, mode="1").convert("L")
    if variant == "sharpen":
        return gray.filter(ImageFilter.UnsharpMask(radius=1.5, percent=175, threshold=3))
    if variant == "invert":
        return ImageOps.invert(gray)
    return gray


def _ocr_one_image(
    image: Image.Image,
    *,
    languages: str,
    psm: int,
    variant: str,
    rotation: int,
    timeout_seconds: int,
) -> tuple[str, float | None, int]:
    import pytesseract
    from pytesseract import Output

    prepared = _prepare_image(image, variant=variant, rotation=rotation)
    data = pytesseract.image_to_data(
        prepared,
        lang=languages,
        config=f"--oem 1 --psm {int(psm)}",
        output_type=Output.DICT,
        timeout=int(timeout_seconds),
    )
    lines: dict[tuple[int, int, int, int], list[str]] = defaultdict(list)
    weighted_confidence = 0.0
    confidence_weight = 0
    for index, raw in enumerate(data.get("text", [])):
        word = str(raw or "").strip()
        try:
            confidence = float(data["conf"][index])
        except (TypeError, ValueError, KeyError):
            confidence = -1.0
        if not word:
            continue
        key = (
            int(data["page_num"][index]),
            int(data["block_num"][index]),
            int(data["par_num"][index]),
            int(data["line_num"][index]),
        )
        lines[key].append(word)
        if confidence >= 0:
            weight = max(len(word), 1)
            weighted_confidence += confidence * weight
            confidence_weight += weight
    text = "\n".join(" ".join(lines[key]) for key in sorted(lines))
    mean_confidence = (
        weighted_confidence / confidence_weight if confidence_weight else None
    )
    return text, mean_confidence, confidence_weight


def _ocr_images(images: list[Image.Image], spec: dict, config: dict) -> tuple[str, float | None]:
    texts: list[str] = []
    confidence_sum = 0.0
    confidence_weight = 0
    for image in images:
        text, confidence, weight = _ocr_one_image(
            image,
            languages=spec["languages"],
            psm=spec["psm"],
            variant=spec["variant"],
            rotation=spec.get("rotation", 0),
            timeout_seconds=int(config["repair"]["tesseract_timeout_seconds"]),
        )
        if text.strip():
            texts.append(text)
        if confidence is not None and weight:
            confidence_sum += confidence * weight
            confidence_weight += weight
    return (
        "\n\n".join(texts),
        confidence_sum / confidence_weight if confidence_weight else None,
    )


def _pdf_page_images(path: Path, page_number: int, dpi: int) -> list[Image.Image]:
    opened = fitz.open(path)
    try:
        if page_number < 1 or page_number > opened.page_count:
            raise IndexError(f"PDF page {page_number} outside 1..{opened.page_count}")
        page = opened[page_number - 1]
        matrix = fitz.Matrix(dpi / 72.0, dpi / 72.0)
        pixmap = page.get_pixmap(matrix=matrix, alpha=False)
        return [Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)]
    finally:
        opened.close()


def _raster_page_images(path: Path, page_number: int) -> list[Image.Image]:
    with Image.open(path) as opened:
        frames = [frame.convert("RGB") for frame in ImageSequence.Iterator(opened)]
    if page_number < 1 or page_number > len(frames):
        raise IndexError(f"Image page {page_number} outside 1..{len(frames)}")
    return [frames[page_number - 1]]


def _pptx_page_images(path: Path, page_number: int) -> tuple[list[Image.Image], str]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(path)
    if page_number < 1 or page_number > len(presentation.slides):
        raise IndexError(f"Slide {page_number} outside 1..{len(presentation.slides)}")
    slide = presentation.slides[page_number - 1]
    native_text: list[str] = []
    images: list[Image.Image] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            native_text.append(shape.text.strip())
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                with Image.open(io.BytesIO(shape.image.blob)) as opened:
                    images.append(opened.convert("RGB"))
            except Exception:
                continue
    return images, "\n".join(native_text)


def _load_source_page(path: Path, page: dict, config: dict) -> tuple[list[Image.Image], str]:
    suffix = path.suffix.casefold()
    page_number = int(page["page_number_1_based"])
    if suffix == ".pdf" or page.get("source_mime_type") == "application/pdf":
        return _pdf_page_images(path, page_number, int(config["repair"]["render_dpi"])), ""
    if suffix in {".jpg", ".jpeg", ".png", ".tif", ".tiff"}:
        return _raster_page_images(path, page_number), ""
    if suffix == ".pptx":
        return _pptx_page_images(path, page_number)
    raise RuntimeError(f"Targeted repair does not support {suffix or page.get('source_mime_type')}")


def _attempt_specs(original: dict, has_native_text: bool, config: dict) -> list[dict]:
    repair = config["repair"]
    arabic_evidence = int(original.get("arabic_character_count") or 0) >= 8
    primary = repair["arabic_languages"] if arabic_evidence else repair["latin_languages"]
    specs: list[dict] = []
    if has_native_text:
        specs.append({"kind": "native", "name": "native_slide_text"})
    specs.extend(
        [
            {"kind": "ocr", "name": "primary_psm3_auto", "languages": primary, "psm": 3, "variant": "autocontrast", "rotation": 0},
            {"kind": "ocr", "name": "primary_psm6_auto", "languages": primary, "psm": 6, "variant": "autocontrast", "rotation": 0},
            {"kind": "ocr", "name": "primary_psm6_threshold", "languages": primary, "psm": 6, "variant": "threshold", "rotation": 0},
            {"kind": "ocr", "name": "primary_psm11_sharpen", "languages": primary, "psm": 11, "variant": "sharpen", "rotation": 0},
        ]
    )
    if original.get("qa_status") == "FAILED":
        specs.extend(
            [
                {"kind": "ocr", "name": "mixed_psm6_auto", "languages": repair["mixed_languages"], "psm": 6, "variant": "autocontrast", "rotation": 0},
                {"kind": "ocr", "name": "primary_psm6_rot90", "languages": primary, "psm": 6, "variant": "autocontrast", "rotation": 90},
                {"kind": "ocr", "name": "primary_psm6_rot270", "languages": primary, "psm": 6, "variant": "autocontrast", "rotation": 270},
            ]
        )
    elif arabic_evidence:
        specs.append(
            {"kind": "ocr", "name": "mixed_psm6_auto", "languages": repair["mixed_languages"], "psm": 6, "variant": "autocontrast", "rotation": 0}
        )
    return specs[: int(repair["max_attempts_per_page"])]


def _original_attempt(original: dict, config: dict) -> dict:
    text = normalize_extracted_text(original.get("text_raw", ""))
    metrics = text_metrics(text)
    confidence = original.get("ocr_confidence")
    if pd.isna(confidence):
        confidence = None
    status, reasons = _classification(metrics, confidence, config)
    return {
        "attempt_order": 0,
        "attempt_name": "original_phase3",
        "attempt_kind": "original",
        "ocr_languages": original.get("ocr_languages", ""),
        "psm": None,
        "preprocessing": "phase3_original",
        "rotation": 0,
        "ocr_confidence": confidence,
        **metrics,
        "qa_status": status,
        "qa_reasons": reasons,
        "text_raw": text,
        "error_type": "",
        "error_message": "",
    }


def repair_page(source_path: Path, original: dict, config: dict) -> tuple[dict, list[dict], list[Image.Image]]:
    images, native_text = _load_source_page(source_path, original, config)
    if not images and not native_text.strip():
        raise RuntimeError("No renderable page pixels or native slide text")
    attempts = [_original_attempt(original, config)]
    specs = _attempt_specs(original, bool(native_text.strip()), config)
    for order, spec in enumerate(specs, start=1):
        try:
            if spec["kind"] == "native":
                text = native_text
                confidence = 100.0
                languages = "native"
            else:
                text, confidence = _ocr_images(images, spec, config)
                languages = spec["languages"]
            text = normalize_extracted_text(text)
            metrics = text_metrics(text)
            status, reasons = _classification(metrics, confidence, config)
            attempt = {
                "attempt_order": order,
                "attempt_name": spec["name"],
                "attempt_kind": spec["kind"],
                "ocr_languages": languages,
                "psm": spec.get("psm"),
                "preprocessing": spec.get("variant", "native"),
                "rotation": spec.get("rotation", 0),
                "ocr_confidence": confidence,
                **metrics,
                "qa_status": status,
                "qa_reasons": reasons,
                "text_raw": text,
                "error_type": "",
                "error_message": "",
            }
        except Exception as exc:
            attempt = {
                "attempt_order": order,
                "attempt_name": spec["name"],
                "attempt_kind": spec["kind"],
                "ocr_languages": spec.get("languages", ""),
                "psm": spec.get("psm"),
                "preprocessing": spec.get("variant", ""),
                "rotation": spec.get("rotation", 0),
                "ocr_confidence": None,
                **text_metrics(""),
                "qa_status": "FAILED",
                "qa_reasons": ["repair_attempt_error"],
                "text_raw": "",
                "error_type": type(exc).__name__,
                "error_message": str(exc)[:300],
            }
        attempts.append(attempt)
        if (
            attempt["qa_status"] == "PASS"
            and float(attempt.get("ocr_confidence") or 0.0)
            >= float(config["repair"]["early_stop_confidence"])
        ):
            break
    selected = max(attempts, key=_attempt_score)
    selected_text = normalize_extracted_text(selected["text_raw"])
    redacted, redaction_counts = redact_text(selected_text, config)
    if selected["attempt_kind"] == "native":
        selected_method, selected_engine = "targeted_native_slide", ""
    elif selected["attempt_kind"] == "ocr":
        selected_method, selected_engine = "targeted_tesseract_ocr", "tesseract"
    else:
        selected_method = original["extraction_method"]
        selected_engine = original.get("ocr_engine", "")
    page_record = dict(original)
    page_record.update(
        {
            "pipeline_version": config["pipeline_version"],
            "extraction_method": selected_method,
            "ocr_engine": selected_engine,
            "ocr_languages": selected["ocr_languages"],
            "ocr_confidence": selected["ocr_confidence"],
            "ink_ratio": max((image_ink_ratio(image) for image in images), default=None),
            "is_blank": False,
            "qa_status": selected["qa_status"],
            "qa_reasons": json.dumps(selected["qa_reasons"], ensure_ascii=False),
            "redaction_counts": json.dumps(redaction_counts, ensure_ascii=False, sort_keys=True),
            "text_raw": selected_text if config["security"].get("persist_raw_text", True) else "",
            "text_redacted": redacted,
            "repair_selected_attempt": selected["attempt_name"],
            "repair_original_qa_status": original["qa_status"],
            "repair_original_ocr_confidence": original.get("ocr_confidence"),
            "repair_original_text_sha256": _text_sha256(normalize_extracted_text(original.get("text_raw", ""))),
            "repair_text_changed": _text_sha256(selected_text) != _text_sha256(normalize_extracted_text(original.get("text_raw", ""))),
            "retrieval_eligible": selected["qa_status"] == "PASS",
            **{key: selected[key] for key in text_metrics("")},
        }
    )
    for attempt in attempts:
        attempt.update(
            {
                "document_id": original["document_id"],
                "source_file_name": original["source_file_name"],
                "page_number_1_based": int(original["page_number_1_based"]),
                "source_sha256": original["source_sha256"],
                "text_sha256": _text_sha256(attempt["text_raw"]),
                "selected": attempt["attempt_name"] == selected["attempt_name"],
            }
        )
    return page_record, attempts, images


def _checkpoint_valid(path: Path, source_sha256: str, config_sha256: str) -> bool:
    if not path.exists():
        return False
    try:
        value = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return False
    return (
        value.get("status") == "COMPLETE"
        and value.get("source_sha256") == source_sha256
        and value.get("config_sha256") == config_sha256
    )


def _preview_montage(images: list[Image.Image], path: Path, max_dimension: int) -> None:
    prepared: list[Image.Image] = []
    for image in images:
        item = ImageOps.exif_transpose(image.copy()).convert("RGB")
        item.thumbnail((max_dimension, max_dimension), Image.Resampling.LANCZOS)
        prepared.append(item)
    if not prepared:
        return
    width = max(image.width for image in prepared)
    height = sum(image.height for image in prepared) + 12 * (len(prepared) - 1)
    canvas = Image.new("RGB", (width, height), "white")
    y = 0
    for image in prepared:
        canvas.paste(image, ((width - image.width) // 2, y))
        y += image.height + 12
    path.parent.mkdir(parents=True, exist_ok=True)
    canvas.save(path, format="JPEG", quality=88, optimize=True)


def _hash_outputs(run_root: Path, paths: Iterable[Path]) -> Path:
    lines = []
    for path in sorted(paths, key=lambda item: str(item.relative_to(run_root))):
        lines.append(f"{sha256_file(path)}  {path.relative_to(run_root)}\n")
    sums = run_root / "SHA256SUMS.txt"
    sums.write_text("".join(lines), encoding="utf-8")
    return sums


def _build_report(summary: dict, unresolved_names: list[str]) -> str:
    names = "\n".join(f"- {name}" for name in unresolved_names) or "- None"
    return f"""# Phase 3.1 — Targeted QA repair

**Technical status:** {summary['status']}
**Quality gate:** {summary['qa_gate']}
**Completed (UTC):** {summary['completed_at_utc']}

## Results

- Pages targeted: **{summary['pages_targeted']}**
- Pages repaired to PASS: **{summary['pages_repaired_to_pass']}**
- Pages still REVIEW: **{summary['pages_still_review']}**
- Pages still FAILED: **{summary['pages_still_failed']}**
- Curated pages: **{summary['curated_pages_total']}**
- Retrieval-eligible pages: **{summary['retrieval_eligible_pages']}**
- Repair processing failures: **{summary['repair_processing_failures']}**

## Unresolved source documents

{names}

## Downstream decision

Only rows with `retrieval_eligible=true` may enter Phase 4 retrieval indexes.
Unresolved pages remain preserved with provenance and are excluded from claims.
No source file or Phase 3 output was modified, and no external LLM was called.
"""


def _verify_inputs(project_root: Path, config: dict) -> tuple[Path, Path, pd.DataFrame, pd.DataFrame]:
    snapshot_id = config["input"]["snapshot_id"]
    snapshot_root = project_root / "data" / "snapshots" / snapshot_id
    phase3_root = (
        project_root
        / config["output"]["root"]
        / snapshot_id
        / config["input"]["phase3_run_name"]
    )
    phase3_manifest = verify_phase3(phase3_root)
    if sha256_file(phase3_root / "PHASE_3_MANIFEST.json") != config["input"]["expected_phase3_manifest_sha256"]:
        raise AssertionError("Pinned Phase 3 manifest hash mismatch")
    if sha256_file(phase3_root / "pages.parquet") != config["input"]["expected_phase3_pages_sha256"]:
        raise AssertionError("Pinned Phase 3 pages hash mismatch")
    pages = pd.read_parquet(phase3_root / "pages.parquet")
    expected_total = int(config["input"]["expected_total_pages"])
    if len(pages) != expected_total or int(phase3_manifest["pages_total"]) != expected_total:
        raise AssertionError("Pinned Phase 3 page count mismatch")
    queue = pages.loc[pages["qa_status"].isin(["REVIEW", "FAILED"])].copy()
    if len(queue) != int(config["input"]["expected_target_pages"]):
        raise AssertionError("Targeted QA page count mismatch")
    counts = queue["qa_status"].value_counts().to_dict()
    if counts.get("REVIEW", 0) != int(config["input"]["expected_review_pages"]):
        raise AssertionError("Pinned REVIEW count mismatch")
    if counts.get("FAILED", 0) != int(config["input"]["expected_failed_pages"]):
        raise AssertionError("Pinned FAILED count mismatch")
    if queue.duplicated(KEY_COLUMNS).any():
        raise AssertionError("Duplicate page keys in Phase 3 QA queue")
    return snapshot_root, phase3_root, pages, queue


def run_phase3_1(project_root: Path, config_path: Path, progress=print) -> dict:
    config = load_phase3_1_config(config_path)
    snapshot_root, phase3_root, pages, queue = _verify_inputs(project_root, config)
    run_root = (
        project_root
        / config["output"]["root"]
        / config["input"]["snapshot_id"]
        / config["output"]["run_name"]
    )
    success_path = run_root / config["output"]["success_marker"]
    if success_path.exists():
        verified = verify_phase3_1(run_root)
        progress("Existing completed Phase 3.1 run verified; no pages were reprocessed.")
        return verified | {"run_root": str(run_root), "resumed": True}

    checkpoints = run_root / "checkpoints"
    previews = run_root / "review_previews"
    reports = run_root / "reports"
    for folder in (checkpoints, previews, reports):
        folder.mkdir(parents=True, exist_ok=True)
    config_hash = sha256_file(config_path)
    processing_errors: list[dict] = []

    queue_records = json.loads(queue.to_json(orient="records", force_ascii=False))
    for index, original in enumerate(queue_records, start=1):
        document_id = original["document_id"]
        page_number = int(original["page_number_1_based"])
        checkpoint = checkpoints / f"{document_id}__p{page_number:04d}.json"
        if _checkpoint_valid(checkpoint, original["source_sha256"], config_hash):
            progress(f"[{index}/{len(queue_records)}] Resume: {original['source_file_name']} p.{page_number}")
            continue
        progress(f"[{index}/{len(queue_records)}] Repairing: {original['source_file_name']} p.{page_number}")
        source_path = snapshot_root / original["source_relative_path"]
        try:
            if sha256_file(source_path) != original["source_sha256"]:
                raise AssertionError("Source hash mismatch before targeted repair")
            page_record, attempts, images = repair_page(source_path, original, config)
            payload = {
                "status": "COMPLETE",
                "config_sha256": config_hash,
                "source_sha256": original["source_sha256"],
                "page_record": page_record,
                "attempts": attempts,
            }
            _atomic_json(checkpoint, payload)
            if page_record["qa_status"] != "PASS":
                _preview_montage(
                    images,
                    previews / f"{document_id}__p{page_number:04d}.jpg",
                    int(config["repair"]["preview_max_dimension"]),
                )
        except Exception as exc:
            processing_errors.append(
                {
                    "document_id": document_id,
                    "source_file_name": original["source_file_name"],
                    "page_number_1_based": page_number,
                    "error_type": type(exc).__name__,
                    "error_message": str(exc)[:500],
                }
            )

    completed_payloads = []
    for path in sorted(checkpoints.glob("*.json")):
        value = json.loads(path.read_text(encoding="utf-8"))
        if value.get("status") == "COMPLETE" and value.get("config_sha256") == config_hash:
            completed_payloads.append(value)
    expected_targets = int(config["input"]["expected_target_pages"])
    if len(completed_payloads) != expected_targets:
        error_path = run_root / "repair_processing_errors.jsonl"
        error_path.write_text(
            "".join(json.dumps(row, ensure_ascii=False) + "\n" for row in processing_errors),
            encoding="utf-8",
        )
        _atomic_json(
            run_root / "_PARTIAL.json",
            {
                "status": "PARTIAL_RETRY_REQUIRED",
                "pages_expected": expected_targets,
                "pages_completed": len(completed_payloads),
                "created_at_utc": _now(),
            },
        )
        raise RuntimeError(
            f"Targeted repair completed {len(completed_payloads)}/{expected_targets} pages. Rerun to retry only errors."
        )

    repaired_rows = [value["page_record"] for value in completed_payloads]
    attempt_rows = [attempt for value in completed_payloads for attempt in value["attempts"]]
    repaired = pd.DataFrame(repaired_rows)
    attempts = pd.DataFrame(attempt_rows)
    repaired_by_key = {
        (row["document_id"], int(row["page_number_1_based"])): row
        for row in repaired_rows
    }

    curated = pages.copy()
    curated["phase3_1_repair_status"] = "NOT_REQUIRED"
    curated["repair_selected_attempt"] = ""
    curated["repair_original_qa_status"] = ""
    curated["repair_original_ocr_confidence"] = pd.NA
    curated["repair_original_text_sha256"] = ""
    curated["repair_text_changed"] = False
    curated["retrieval_eligible"] = curated["qa_status"].eq("PASS")
    for row_index, row in curated.iterrows():
        key = (row["document_id"], int(row["page_number_1_based"]))
        repaired_row = repaired_by_key.get(key)
        if repaired_row is None:
            continue
        for column in FINAL_TEXT_COLUMNS:
            curated.at[row_index, column] = repaired_row[column]
        for column in (
            "repair_selected_attempt",
            "repair_original_qa_status",
            "repair_original_ocr_confidence",
            "repair_original_text_sha256",
            "repair_text_changed",
            "retrieval_eligible",
        ):
            curated.at[row_index, column] = repaired_row[column]
        curated.at[row_index, "pipeline_version"] = config["pipeline_version"]
        curated.at[row_index, "phase3_1_repair_status"] = (
            "REPAIRED_PASS" if repaired_row["qa_status"] == "PASS" else "UNRESOLVED_EXCLUDED"
        )

    repaired_path = run_root / "repaired_pages.parquet"
    attempts_path = run_root / "repair_attempts.parquet"
    curated_path = run_root / "pages_curated.parquet"
    repaired.to_parquet(repaired_path, index=False)
    attempts.to_parquet(attempts_path, index=False)
    curated.to_parquet(curated_path, index=False)

    safe_columns = [
        "document_id",
        "source_file_name",
        "page_number_1_based",
        "repair_original_qa_status",
        "qa_status",
        "repair_selected_attempt",
        "ocr_languages",
        "ocr_confidence",
        "character_count",
        "word_count",
        "line_count",
        "repair_text_changed",
        "retrieval_eligible",
        "source_sha256",
        "source_relative_path",
    ]
    results_csv = run_root / "repair_results.csv"
    repaired[safe_columns].sort_values(KEY_COLUMNS).to_csv(
        results_csv, index=False, encoding="utf-8-sig"
    )
    unresolved = repaired.loc[~repaired["qa_status"].eq("PASS")].copy()
    unresolved_csv = run_root / "unresolved_review_queue.csv"
    unresolved[safe_columns].sort_values(KEY_COLUMNS).to_csv(
        unresolved_csv, index=False, encoding="utf-8-sig"
    )
    errors_path = run_root / "repair_processing_errors.jsonl"
    errors_path.write_text("", encoding="utf-8")

    counts = Counter(repaired["qa_status"].tolist())
    unresolved_names = sorted(unresolved["source_file_name"].dropna().unique().tolist())
    summary = {
        "schema_version": 1,
        "phase": 3.1,
        "pipeline_version": config["pipeline_version"],
        "snapshot_id": config["input"]["snapshot_id"],
        "phase3_manifest_sha256": config["input"]["expected_phase3_manifest_sha256"],
        "phase3_pages_sha256": config["input"]["expected_phase3_pages_sha256"],
        "status": "PASS",
        "qa_gate": "PASS" if counts.get("REVIEW", 0) == 0 and counts.get("FAILED", 0) == 0 else "REVIEW_REQUIRED",
        "pages_targeted": int(len(repaired)),
        "pages_repaired_to_pass": int(counts.get("PASS", 0)),
        "pages_still_review": int(counts.get("REVIEW", 0)),
        "pages_still_failed": int(counts.get("FAILED", 0)),
        "curated_pages_total": int(len(curated)),
        "retrieval_eligible_pages": int(curated["retrieval_eligible"].sum()),
        "retrieval_excluded_pages": int((~curated["retrieval_eligible"]).sum()),
        "repair_processing_failures": 0,
        "unresolved_documents": int(unresolved["document_id"].nunique()),
        "source_snapshot_mutation_calls": 0,
        "phase3_output_mutation_calls": 0,
        "external_llm_calls": 0,
        "completed_at_utc": _now(),
    }
    manifest_path = run_root / "PHASE_3_1_MANIFEST.json"
    _atomic_json(manifest_path, summary)
    report_path = reports / "PHASE_3_1_REPAIR_REPORT.md"
    report_path.write_text(_build_report(summary, unresolved_names), encoding="utf-8")
    final_paths = [
        repaired_path,
        attempts_path,
        curated_path,
        results_csv,
        unresolved_csv,
        errors_path,
        manifest_path,
        report_path,
        *sorted(previews.glob("*.jpg")),
    ]
    sums_path = _hash_outputs(run_root, final_paths)
    success = {
        "status": "COMPLETE_REPRODUCIBLE_TARGETED_REPAIR",
        "snapshot_id": config["input"]["snapshot_id"],
        "pipeline_version": config["pipeline_version"],
        "manifest_sha256": sha256_file(manifest_path),
        "sha256sums_sha256": sha256_file(sums_path),
        "created_at_utc": _now(),
    }
    partial = run_root / "_PARTIAL.json"
    if partial.exists():
        partial.unlink()
    _atomic_json(success_path, success)
    return summary | {"run_root": str(run_root), "resumed": False}


def verify_phase3_1(run_root: Path) -> dict:
    success_path = run_root / "_SUCCESS.json"
    manifest_path = run_root / "PHASE_3_1_MANIFEST.json"
    sums_path = run_root / "SHA256SUMS.txt"
    if not all(path.exists() for path in (success_path, manifest_path, sums_path)):
        raise AssertionError("Incomplete Phase 3.1 output")
    success = json.loads(success_path.read_text(encoding="utf-8"))
    if sha256_file(manifest_path) != success["manifest_sha256"]:
        raise AssertionError("Phase 3.1 manifest hash mismatch")
    if sha256_file(sums_path) != success["sha256sums_sha256"]:
        raise AssertionError("Phase 3.1 checksum-list hash mismatch")
    for line in sums_path.read_text(encoding="utf-8").splitlines():
        expected, relative = line.split("  ", 1)
        if sha256_file(run_root / relative) != expected:
            raise AssertionError(f"Phase 3.1 output hash mismatch: {relative}")
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    if manifest.get("status") != "PASS":
        raise AssertionError("Phase 3.1 technical status is not PASS")
    if manifest.get("pages_targeted") != 21 or manifest.get("curated_pages_total") != 408:
        raise AssertionError("Phase 3.1 pinned page counts changed")
    if any(
        manifest.get(field) != 0
        for field in (
            "repair_processing_failures",
            "source_snapshot_mutation_calls",
            "phase3_output_mutation_calls",
            "external_llm_calls",
        )
    ):
        raise AssertionError("Phase 3.1 safety assertion failed")
    return manifest
