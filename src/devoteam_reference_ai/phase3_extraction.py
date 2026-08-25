from __future__ import annotations

import shutil
import subprocess
import tempfile
from collections import Counter, defaultdict
from pathlib import Path

import fitz
from PIL import Image, ImageOps, ImageSequence

from .phase3_quality import (
    image_ink_ratio,
    is_usable_digital_text,
    normalize_extracted_text,
    page_quality,
    text_metrics,
)
from .phase3_security import redact_text


PDF_MIME = "application/pdf"
IMAGE_MIMES = {"image/jpeg", "image/png", "image/tiff"}
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _page_record(
    *,
    document: dict,
    page_number: int,
    text: str,
    method: str,
    ocr_confidence: float | None,
    ink_ratio: float | None,
    config: dict,
) -> dict:
    normalized = normalize_extracted_text(text)
    metrics = text_metrics(normalized)
    blank = ink_ratio is not None and ink_ratio < float(
        config["extraction"]["blank_ink_ratio_threshold"]
    )
    status, reasons = page_quality(method, metrics, ocr_confidence, blank, config)
    redacted, redaction_counts = redact_text(normalized, config)
    return {
        "snapshot_id": config["input"]["snapshot_id"],
        "pipeline_version": config["pipeline_version"],
        "document_id": document["file_id"],
        "source_file_name": document["name"],
        "source_mime_type": document["mimeType"],
        "source_sha256": document["local_sha256"],
        "source_relative_path": document["local_relative_path"],
        "page_number_1_based": page_number,
        "extraction_method": method,
        "ocr_engine": "tesseract" if method == "tesseract_ocr" else "",
        "ocr_languages": config["extraction"]["ocr_languages"] if method == "tesseract_ocr" else "",
        "ocr_confidence": ocr_confidence,
        "ink_ratio": ink_ratio,
        **metrics,
        "is_blank": blank,
        "qa_status": status,
        "qa_reasons": reasons,
        "security_classification": config["security"]["classification"],
        "redaction_counts": redaction_counts,
        "text_raw": normalized if config["security"].get("persist_raw_text", True) else "",
        "text_redacted": redacted,
    }


def _render_page(page, dpi: int) -> Image.Image:
    matrix = fitz.Matrix(dpi / 72.0, dpi / 72.0)
    pixmap = page.get_pixmap(matrix=matrix, alpha=False)
    return Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)


def _ocr_image(image: Image.Image, config: dict) -> tuple[str, float | None]:
    import pytesseract
    from pytesseract import Output

    prepared = ImageOps.autocontrast(ImageOps.grayscale(ImageOps.exif_transpose(image)))
    data = pytesseract.image_to_data(
        prepared,
        lang=config["extraction"]["ocr_languages"],
        config=f"--oem 1 --psm {int(config['extraction']['tesseract_psm'])}",
        output_type=Output.DICT,
    )
    lines: dict[tuple[int, int, int, int], list[str]] = defaultdict(list)
    confidences: list[float] = []
    for index, raw in enumerate(data.get("text", [])):
        word = str(raw or "").strip()
        try:
            confidence = float(data["conf"][index])
        except (TypeError, ValueError, KeyError):
            confidence = -1.0
        if word:
            key = (
                int(data["page_num"][index]),
                int(data["block_num"][index]),
                int(data["par_num"][index]),
                int(data["line_num"][index]),
            )
            lines[key].append(word)
            if confidence >= 0:
                confidences.append(confidence)
    text = "\n".join(" ".join(lines[key]) for key in sorted(lines))
    mean_confidence = sum(confidences) / len(confidences) if confidences else None
    return text, mean_confidence


def _extract_pdf(pdf_path: Path, document: dict, config: dict) -> list[dict]:
    records: list[dict] = []
    pdf = fitz.open(pdf_path)
    try:
        for page_index, page in enumerate(pdf):
            digital = normalize_extracted_text(page.get_text("text"))
            if is_usable_digital_text(digital, config):
                records.append(
                    _page_record(
                        document=document,
                        page_number=page_index + 1,
                        text=digital,
                        method="digital_pdf",
                        ocr_confidence=None,
                        ink_ratio=None,
                        config=config,
                    )
                )
                continue
            image = _render_page(page, int(config["extraction"]["pdf_render_dpi"]))
            ink_ratio = image_ink_ratio(image)
            if ink_ratio < float(config["extraction"]["blank_ink_ratio_threshold"]):
                text, confidence = "", None
            else:
                text, confidence = _ocr_image(image, config)
            records.append(
                _page_record(
                    document=document,
                    page_number=page_index + 1,
                    text=text,
                    method="tesseract_ocr",
                    ocr_confidence=confidence,
                    ink_ratio=ink_ratio,
                    config=config,
                )
            )
    finally:
        pdf.close()
    return records


def _extract_images(path: Path, document: dict, config: dict) -> list[dict]:
    records: list[dict] = []
    with Image.open(path) as opened:
        for page_number, frame in enumerate(ImageSequence.Iterator(opened), start=1):
            image = frame.convert("RGB")
            ink_ratio = image_ink_ratio(image)
            if ink_ratio < float(config["extraction"]["blank_ink_ratio_threshold"]):
                text, confidence = "", None
            else:
                text, confidence = _ocr_image(image, config)
            records.append(
                _page_record(
                    document=document,
                    page_number=page_number,
                    text=text,
                    method="tesseract_ocr",
                    ocr_confidence=confidence,
                    ink_ratio=ink_ratio,
                    config=config,
                )
            )
    return records


def _convert_office_to_pdf(source: Path, output_dir: Path, timeout: int) -> Path:
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if not executable:
        raise RuntimeError("LibreOffice is not installed")
    profile = output_dir / "lo_profile"
    command = [
        executable,
        "--headless",
        f"-env:UserInstallation=file://{profile}",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(source),
    ]
    result = subprocess.run(command, capture_output=True, text=True, timeout=timeout)
    candidates = list(output_dir.glob("*.pdf"))
    if result.returncode != 0 or len(candidates) != 1:
        raise RuntimeError(f"LibreOffice conversion failed: {result.stderr[-300:]}")
    return candidates[0]


def _native_office_fallback(path: Path, document: dict, config: dict) -> list[dict]:
    if document["mimeType"] == DOCX_MIME or path.suffix.casefold() == ".docx":
        from docx import Document

        opened = Document(path)
        blocks = [paragraph.text for paragraph in opened.paragraphs if paragraph.text.strip()]
        for table in opened.tables:
            blocks.extend(
                " | ".join(cell.text.strip() for cell in row.cells)
                for row in table.rows
                if any(cell.text.strip() for cell in row.cells)
            )
        text = "\n".join(blocks)
        return [
            _page_record(
                document=document,
                page_number=1,
                text=text,
                method="docx_native_fallback",
                ocr_confidence=None,
                ink_ratio=None,
                config=config,
            )
        ]
    if document["mimeType"] == PPTX_MIME or path.suffix.casefold() == ".pptx":
        from pptx import Presentation

        presentation = Presentation(path)
        records = []
        for page_number, slide in enumerate(presentation.slides, start=1):
            text = "\n".join(
                shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()
            )
            records.append(
                _page_record(
                    document=document,
                    page_number=page_number,
                    text=text,
                    method="pptx_native_fallback",
                    ocr_confidence=None,
                    ink_ratio=None,
                    config=config,
                )
            )
        return records
    raise RuntimeError(f"Unsupported office file: {path.suffix}")


def process_document(path: Path, document: dict, config: dict) -> tuple[dict, list[dict]]:
    try:
        mime_type = document["mimeType"]
        if mime_type == PDF_MIME or path.suffix.casefold() == ".pdf":
            pages = _extract_pdf(path, document, config)
        elif mime_type in IMAGE_MIMES or path.suffix.casefold() in {".jpg", ".jpeg", ".png", ".tif", ".tiff"}:
            pages = _extract_images(path, document, config)
        elif mime_type in {DOCX_MIME, PPTX_MIME} or path.suffix.casefold() in {".docx", ".pptx"}:
            with tempfile.TemporaryDirectory(prefix=f"phase3_{document['file_id']}_") as temporary:
                try:
                    converted = _convert_office_to_pdf(
                        path,
                        Path(temporary),
                        int(config["extraction"]["office_conversion_timeout_seconds"]),
                    )
                    pages = _extract_pdf(converted, document, config)
                except Exception:
                    pages = _native_office_fallback(path, document, config)
        else:
            raise RuntimeError(f"Unsupported MIME type: {mime_type}")
        method_counts = Counter(page["extraction_method"] for page in pages)
        qa_counts = Counter(page["qa_status"] for page in pages)
        summary = {
            "snapshot_id": config["input"]["snapshot_id"],
            "pipeline_version": config["pipeline_version"],
            "document_id": document["file_id"],
            "source_file_name": document["name"],
            "source_mime_type": document["mimeType"],
            "source_sha256": document["local_sha256"],
            "source_relative_path": document["local_relative_path"],
            "security_classification": config["security"]["classification"],
            "status": "COMPLETE",
            "page_count": len(pages),
            "method_counts": dict(method_counts),
            "qa_counts": dict(qa_counts),
            "error_type": "",
            "error_message": "",
        }
        return summary, pages
    except Exception as exc:
        return (
            {
                "snapshot_id": config["input"]["snapshot_id"],
                "pipeline_version": config["pipeline_version"],
                "document_id": document["file_id"],
                "source_file_name": document["name"],
                "source_mime_type": document["mimeType"],
                "source_sha256": document["local_sha256"],
                "source_relative_path": document["local_relative_path"],
                "security_classification": config["security"]["classification"],
                "status": "FAILED",
                "page_count": 0,
                "method_counts": {},
                "qa_counts": {},
                "error_type": type(exc).__name__,
                "error_message": str(exc)[:500],
            },
            [],
        )
